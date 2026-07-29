import os
import sys
import uuid
import chromadb
from dotenv import load_dotenv
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader

load_dotenv()

MISTRAL_API_KEY = os.getenv("MISTRAL_API_KEY")
_embedding_model = None
_llm = None


def _require_api_key():
    if not MISTRAL_API_KEY:
        raise RuntimeError(
            "MISTRAL_API_KEY is not set. Copy .env.example to .env and add your key "
            "from https://console.mistral.ai"
        )


def get_embedding_model():
    global _embedding_model
    _require_api_key()
    if _embedding_model is None:
        from langchain_mistralai import MistralAIEmbeddings

        _embedding_model = MistralAIEmbeddings(api_key=MISTRAL_API_KEY)
    return _embedding_model


def get_llm():
    global _llm
    _require_api_key()
    if _llm is None:
        from langchain_mistralai import ChatMistralAI

        _llm = ChatMistralAI(api_key=MISTRAL_API_KEY, model="mistral-small-latest")
    return _llm


def get_chroma_client():
    return chromadb.PersistentClient(path="./chroma_db")


def _pptx_available():
    try:
        import pptx  # noqa: F401

        return True
    except ImportError:
        return False


def _ocr_available():
    try:
        import pytesseract  # noqa: F401
        from PIL import Image  # noqa: F401

        return True
    except ImportError:
        return False


def get_supported_formats():
    formats = ["pdf", "docx"]
    if _pptx_available():
        formats.append("pptx")
    if _ocr_available() and sys.version_info < (3, 13):
        formats.extend(["png", "jpg", "jpeg"])
    return formats


def _load_pptx(file_path):
    from pptx import Presentation

    presentation = Presentation(file_path)
    texts = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            texts.append(f"Slide {slide_idx}:\n" + "\n".join(parts))
    if not texts:
        return []
    return [Document(page_content="\n\n".join(texts), metadata={"source": file_path})]


def _load_image(file_path):
    import pytesseract
    from PIL import Image

    text = pytesseract.image_to_string(Image.open(file_path)).strip()
    if not text:
        return []
    return [Document(page_content=text, metadata={"source": file_path})]


def extract_text(file_path):
    ext = os.path.splitext(file_path)[-1].lower()

    if ext == ".pdf":
        return PyPDFLoader(file_path).load()
    if ext == ".docx":
        return Docx2txtLoader(file_path).load()
    if ext == ".pptx":
        if not _pptx_available():
            raise ValueError("PPTX support requires python-pptx. Install with: pip install python-pptx")
        return _load_pptx(file_path)
    if ext in {".png", ".jpg", ".jpeg"}:
        if not _ocr_available():
            raise ValueError(
                "Image support requires pillow and pytesseract, plus the Tesseract OCR system binary."
            )
        return _load_image(file_path)

    raise ValueError(f"Unsupported file type: {ext}")


def process_and_embed(file_path):
    docs = extract_text(file_path)
    if not docs:
        return []

    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    chunks = splitter.split_documents(docs)
    texts = [c.page_content.strip() for c in chunks if c.page_content and c.page_content.strip()]
    if not texts:
        return []

    vectors = get_embedding_model().embed_documents(texts)
    return [
        {
            "id": str(uuid.uuid4()),
            "text": texts[i],
            "embedding": vectors[i],
        }
        for i in range(len(texts))
    ]


def save_to_chromadb(chunks, client, collection_name):
    try:
        collection = client.get_collection(collection_name)
    except Exception:
        collection = client.create_collection(collection_name)

    collection.add(
        ids=[c["id"] for c in chunks],
        documents=[c["text"] for c in chunks],
        embeddings=[c["embedding"] for c in chunks],
    )


def answer_query(query, client, collection_name, top_k=5):
    try:
        collection = client.get_collection(collection_name)
    except Exception:
        return "No documents uploaded. Upload a PDF or DOCX first."

    if collection.count() == 0:
        return "No documents uploaded. Upload a PDF or DOCX first."

    query_vector = get_embedding_model().embed_query(query)
    results = collection.query(
        query_embeddings=[query_vector],
        n_results=min(top_k, collection.count()),
    )

    docs_lists = results.get("documents") or []
    if not docs_lists or not docs_lists[0]:
        return "No relevant context found."

    context = "\n".join(docs_lists[0])
    prompt = f"""Answer ONLY from the context below.
If the answer is not in the context, say you don't know.

Context:
{context}

Question:
{query}
"""
    return get_llm().invoke(prompt).content
